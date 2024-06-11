import pickle

import requests
import os

import yadisk
from pptx import Presentation
from TGDesignBot.YandexDisk.YaDiskInfo import TemplateInfo
import aspose.slides as slides

ya_disk = yadisk.YaDisk(token=str(os.getenv('YANDEX_DISK_TOKEN')))
# ya_disk = pickle.load(open("YandexDisk/YaDisk.pkl", "rb"))

class SlideInfo:
    def __init__(self, slide_idx: int, tags: str):
        self.template_info = None
        self.slide_idx = slide_idx
        self.tags = tags
        self.idx_list = [slide_idx]

    def add_template_info(self, template_info: TemplateInfo):
        self.template_info = template_info

    def add_id(self, slide_id: int):
        self.idx_list.append(slide_id)

    def add_index(self, index: int):
        self.idx_list.append(index)

    def add_indexes(self, indexes: list):
        self.idx_list += indexes


# Takes path to file in local directory and list of templates (pptx files).
# Install them to a local directory.
def install_templates(path: str, templates: list):
    if not os.path.exists(path):
        os.makedirs(path)

    try:
        for template in templates:
            if os.path.exists(path + template.name):
                continue
            response = requests.get(ya_disk.get_download_link(template.path + '/' + template.name))
            with open(path + template.name, 'wb') as file:
                file.write(response.content)
    except Exception as e:
        for template in templates:
            if os.path.exists(path + template.name):
                os.remove(path + template.file)
        raise Exception("Can't install templates")


# Takes the path to presentation and returns a list of SlideInfo classes. Using 0-indexation.
def get_slides_information(path: str) -> list:
    exists(path)
    slides_info = []
    try:
        with slides.Presentation(path) as presentation:
            for author in presentation.comment_authors:
                for comment in author.comments:
                    slides_info.append(SlideInfo(comment.slide.slide_number - 1,
                                                 ";".join(comment.text.split())))
    except Exception as e:
        print(e)
    finally:
        return slides_info


def exists(path: str):
    if not os.path.exists(path):
        raise Exception("No such file or directory")


# This function delete all comments from presentation placed in path.
def remove_all_comments(path: str):
    exists(path)
    with slides.Presentation(path) as presentation:
        # Deletes all comments from the presentation
        for author in presentation.comment_authors:
            author.comments.clear()
        # Deletes all authors
        presentation.comment_authors.clear()
        presentation.save(path, slides.export.SaveFormat.PPTX)
    remove_all_watermarks(path)


def __remove_all_comments__(presentation: slides.Presentation):
    for author in presentation.comment_authors:
        author.comments.clear()
    # Deletes all authors
    presentation.comment_authors.clear()


def remove_template(path: str):
    if not os.path.exists(path):
        return
    os.remove(path)


# Function takes a path where need to save concatenated slides and SlideInfo struct.
# If template didn't install, install it. And from last slide to first delete slides
# by index if no such index in SlideInfo. At the end remove all comments from output presentation.
def get_template_of_slides(path: str, slide_info: SlideInfo):
    remove_template(path)

    install_templates("./Data/Templates/", [slide_info.template_info])

    with slides.Presentation(f"./Data/Templates/{slide_info.template_info.name}") as source_pres:
        for idx in range(source_pres.slides.length - 1, -1, -1):
            if not (idx in slide_info.idx_list):
                source_pres.slides.remove_at(idx)
        __remove_all_comments__(source_pres)
        source_pres.save(path, slides.export.SaveFormat.PPTX)

    remove_all_watermarks(path)


# This function delete all waste, what was added by asponse from presentation.
def remove_all_watermarks(path: str):
    exists(path)

    presentation = Presentation(path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and len(shape.text_frame.paragraphs) > 1:
                if "Created with Aspose.Slides for Python" in shape.text_frame.paragraphs[1].text:
                    shape.element.getparent().remove(shape.element)
    presentation.save(path)
